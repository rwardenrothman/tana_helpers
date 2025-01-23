from configparser import ConfigParser, NoOptionError
from pathlib import Path

import appdirs
import fitz
from pptx import Presentation
from pptx2md import outputter, g
from pptx2md import parser
from tqdm import tqdm
import typer
import win32com.client

import tana_helpers.TanaAPI as tapi
from tana_helpers.TanaAPI import LargeFileError

SLIDE_TAG = tapi.SuperTag(id='EzfED-izb3-0')
APP_PATH = Path(appdirs.AppDirs('hc2md').user_data_dir)
TABLE_FIELD_PATH = APP_PATH / 'pptx_table_fields.cfg'
TABLE_FIELDS_NODE_ID = '2ZW21NkrZcxq'
FIELD_TAG = tapi.SuperTag(id='SYS_T02')
TABLE_TAG = tapi.SuperTag(id='HyaHwYWkdPXK')


class OutputterData:
    def __init__(self):
        self.title = ''
        self.data: list[tuple[str, int]] = []
        self.index = 0

    def to_nodes(self) -> list[tapi.Node]:
        parent_node = tapi.Node(children=[])
        cur_node_by_level: dict[int, tapi.Node] = {-1: parent_node}
        for cur_data, cur_level in self.data:
            if cur_data == TABLE_TAG.id:
                cur_node = tapi.PlainNode(name='Table', supertags=[TABLE_TAG])
            elif cur_data.endswith('::'):
                cur_node = get_field(cur_data[:-2].strip())
            else:
                cur_node = tapi.PlainNode(name=cur_data)
            cur_node_by_level[cur_level] = cur_node
            try:
                cur_node_by_level[cur_level - 1].children.append(cur_node)
            except AttributeError:
                cur_node_by_level[cur_level - 1].children = [cur_node]
            except KeyError:
                cur_level -= 1
                while cur_level >= 0:
                    cur_level -= 1
                    if cur_level in cur_node_by_level:
                        cur_node_by_level[cur_level].children.append(cur_node)
                        break

        return parent_node.children


class TanaOutputter(outputter.md_outputter):
    def __init__(self):
        # configs
        # disable image extraction
        g.disable_image = True
        # prevent adding html tags with colors
        g.disable_color = True
        # prevent escaping of characters
        g.disable_escaping = True
        g.disable_notes = True
        g.enable_slides = True

        self.cur_slide = OutputterData()
        self.slides: list[OutputterData] = []
        self._next_slide()

    def _next_slide(self):
        self.slides.append(OutputterData())
        self.cur_slide = self.slides[-1]

    def put_header(self):
        pass

    def put_title(self, text, level):
        self.cur_slide.title = text

    def put_list(self, text, level):
        self.cur_slide.data.append((text, level))

    def put_para(self, text):
        if text == "\n---\n":
            self._next_slide()
        else:
            self.cur_slide.data.append((text, 0))

    def put_image(self, path, max_width=None):
        return super().put_image(path, max_width)

    def put_table(self, table):
        headers, *rows = table
        self.cur_slide.data.append((TABLE_TAG.id, 0))
        for i, row in enumerate(rows):
            self.cur_slide.data.append((f'{i+1}', 1))
            for cur_col, cur_val in zip(headers, row):
                if cur_val:
                    self.cur_slide.data.append((f'{cur_col}::', 2))
                    for val_row in cur_val.split('\n'):
                        self.cur_slide.data.append((val_row, 3))
            # self.put_para(' | '.join(row).replace('\n', ' ').strip())

    def get_accent(self, text):
        return f" __{text.strip()}__ "

    def get_strong(self, text):
        return f" **{text.strip()}** "

    def get_colored(self, text, rgb):
        return super().get_colored(text, rgb)

    def get_hyperlink(self, text, url):
        return super().get_hyperlink(text, url)

    def get_escaped(self, text):
        return super().get_escaped(text)

    def write(self, text):
        pass

    def flush(self):
        pass

    def close(self):
        pass


def get_field(field_name: str, make_new: bool = True) -> tapi.FieldNode:
    node_name = field_name.replace(':', '').replace('=', '')
    if TABLE_FIELD_PATH.is_file():
        parser = ConfigParser()
        parser.read(TABLE_FIELD_PATH)
    else:
        parser = ConfigParser()
        with TABLE_FIELD_PATH.open('w') as cfg_file:
            parser.write(cfg_file)

    try:
        node_id = parser.get('DEFAULT', node_name)
    except NoOptionError:
        if not make_new:
            raise

        field_node = tapi.PlainNode(name=field_name, supertags=[FIELD_TAG])
        new_nodes = tapi.Tana(field_node, target_id=TABLE_FIELDS_NODE_ID).submit(True)
        node_id = new_nodes.children[0].nodeId

        parser.set('DEFAULT', node_name, node_id)
        with TABLE_FIELD_PATH.open('w') as cfg_file:
            parser.write(cfg_file)

    return tapi.FieldNode(attributeId=node_id)


class SlideNode(tapi.Node):
    @classmethod
    def new(cls, data: OutputterData, image_name: str, image_bytes: bytes):
        img_field = tapi.FieldNode(attributeId='ldlB7rVr3pTc',
                                   children=[tapi.FileNode.from_bytes(image_bytes, image_name)])

        nodes = data.to_nodes()
        if nodes:
            md_field = tapi.FieldNode(attributeId='hzZFWlILJ6hd', children=nodes)
            slide_node = tapi.Node(name=data.title, supertags=[SLIDE_TAG], children=[img_field, md_field])
        else:
            slide_node = tapi.Node(name=data.title or 'Untitled Slide', supertags=[SLIDE_TAG], children=[img_field])
        return slide_node


def ppt_to_pdf(input_file: Path, output_file: Path, format_type=32):
    powerpoint = win32com.client.DispatchEx("Powerpoint.Application")
    powerpoint.Visible = 1

    if output_file.suffix != '.pdf':
        output_file = output_file.with_suffix('.pdf')
    deck = powerpoint.Presentations.Open(input_file)
    deck.SaveAs(output_file, format_type)  # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()


def main(pptx_path: Path, pdf_path: Path = None, target_node_id: str = None):
    # Get missing values
    if not pdf_path:
        pdf_path = pptx_path.with_suffix('.pdf')
        if not pdf_path.exists():
            print(f'Converting {pptx_path.name} to PDF')
            ppt_to_pdf(pptx_path, pdf_path)

    if target_node_id is None:
        print('Adding a node to your Tana Inbox')
        slides_node = tapi.Tana(tapi.PlainNode(name=f'Slides ({pptx_path.name})')).target_inbox().submit()
        target_node_id = slides_node.children[0].nodeId
        print(f'Complete. See https://app.tana.inc?nodeid={target_node_id}')
    elif target_node_id.startswith('http'):
        target_node_id = target_node_id.split('=', 1)[1]
    api = tapi.Tana().set_target_id(target_node_id)

    print('Opening PPTX')
    prs = Presentation(pptx_path)
    t_out = TanaOutputter()
    print('Extracting Text')
    parser.parse(prs, t_out)
    del prs
    print('Completed Extracting')

    print('Opening PDF')
    pdf = fitz.open(pdf_path.as_posix())

    cur_page: fitz.Page
    for i, (cur_page, slide) in tqdm(enumerate(zip(pdf, t_out.slides)), total=len(pdf)):
        # print(f'Getting Slide {i+1}')
        cur_zoom = 1.0
        while cur_zoom > 0:
            try:
                zoom_matrix = fitz.Matrix(cur_zoom, cur_zoom)
                page_img_bytes = cur_page.get_pixmap(matrix=zoom_matrix).tobytes()
                api.add_children(SlideNode.new(slide, f'slide{i + 1:d}.png', page_img_bytes))
                api.submit()
                break
            except LargeFileError as lfe:
                cur_zoom  *= lfe.MAX_LEN / lfe.encoded_size
        else:
            raise LargeFileError('Could not zoom out enough to make the image uploadable')

    return {
        "statusCode": 200,
        "body": 'complete',
    }


if __name__ == '__main__':
    typer.run(main)
